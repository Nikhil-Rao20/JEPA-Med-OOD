#!/usr/bin/env python3
"""
Generate Modern PowerPoint Presentation for Med_JEPA_ODD Project
Self-Supervised Learning for Medical Imaging OOD Detection
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Get the project root directory
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(PROJECT_ROOT, "Med_JEPA_OOD_Presentation.pptx")

# Path to embedding images
EMBEDDING_DIR = os.path.join(PROJECT_ROOT, "experiments", "cxr_jepa_pilot", "3d_visualization")
ANALYSIS_DIR = os.path.join(PROJECT_ROOT, "experiments", "cxr_jepa_pilot", "embedding_analysis")

# Modern Color Palette - Deep Blue/Teal theme
COLORS = {
    'primary_dark': RGBColor(15, 32, 65),      # Deep navy blue
    'primary': RGBColor(26, 54, 93),           # Navy blue
    'accent': RGBColor(0, 180, 216),           # Bright teal/cyan
    'accent2': RGBColor(144, 224, 239),        # Light cyan
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(240, 248, 255),     # Alice blue
    'text_light': RGBColor(220, 230, 240),
    'text_dark': RGBColor(30, 30, 50),
    'gradient_start': RGBColor(15, 32, 65),
    'gradient_end': RGBColor(44, 82, 130),
}


def add_gradient_background(slide, color1, color2):
    """Add a gradient background to the slide."""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2


def add_solid_background(slide, color):
    """Add a solid color background."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top, height=0.15):
    """Add a horizontal accent bar."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(top), 
        Inches(10), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()


def add_side_accent(slide):
    """Add a vertical accent bar on the left side."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.15), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()


def add_corner_decoration(slide):
    """Add decorative elements in corners."""
    # Top right circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8.5), Inches(-0.5),
        Inches(2), Inches(2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['accent']
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    # Bottom left circle
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-0.8), Inches(6),
        Inches(2.5), Inches(2.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = COLORS['accent2']
    circle2.fill.fore_color.brightness = 0.4
    circle2.line.fill.background()


def add_title_slide(prs, title, subtitle=""):
    """Add a modern title slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark gradient background
    add_solid_background(slide, COLORS['primary_dark'])
    
    # Decorative elements
    add_corner_decoration(slide)
    
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(3.3),
        Inches(2), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(1.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1))
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = COLORS['accent2']
        sp.alignment = PP_ALIGN.LEFT
    
    return slide


def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['primary'])
    
    # Large number or icon area
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7), Inches(1),
        Inches(4), Inches(4)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = COLORS['accent']
    accent_shape.fill.fore_color.brightness = 0.2
    accent_shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    return slide


def add_content_slide(prs, title, bullet_points, two_column=False):
    """Add a modern content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background with accent
    add_solid_background(slide, COLORS['light_gray'])
    add_side_accent(slide)
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    if two_column and len(bullet_points) > 4:
        # Two column layout
        mid = len(bullet_points) // 2
        left_points = bullet_points[:mid]
        right_points = bullet_points[mid:]
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5))
        add_bullet_text(left_box, left_points)
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5))
        add_bullet_text(right_box, right_points)
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        add_bullet_text(content_box, bullet_points)
    
    return slide


def add_bullet_text(textbox, points):
    """Add bullet points to a textbox."""
    ctf = textbox.text_frame
    ctf.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = ctf.paragraphs[0]
        else:
            p = ctf.add_paragraph()
        
        # Check for sub-bullets (starting with spaces)
        if point.startswith("   "):
            p.text = "    • " + point.strip()
            p.font.size = Pt(16)
            p.level = 1
        else:
            p.text = "● " + point
            p.font.size = Pt(18)
        
        p.font.color.rgb = COLORS['text_dark']
        p.space_after = Pt(10)


def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['light_gray'])
    add_side_accent(slide)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Image
    if os.path.exists(image_path):
        img_left = Inches(1.5)
        img_top = Inches(1.5)
        img_width = Inches(7)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
    else:
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(2.5), Inches(6), Inches(3)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(200, 210, 220)
        ptf = placeholder.text_frame
        pp = ptf.paragraphs[0]
        pp.text = f"[Image: {os.path.basename(image_path)}]"
        pp.font.size = Pt(16)
        pp.font.color.rgb = COLORS['text_dark']
        pp.alignment = PP_ALIGN.CENTER
    
    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        ctf = cap_box.text_frame
        cp = ctf.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(12)
        cp.font.italic = True
        cp.font.color.rgb = COLORS['text_dark']
        cp.alignment = PP_ALIGN.CENTER
    
    return slide


def add_multi_image_slide(prs, title, images_with_captions):
    """Add a slide with multiple images."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['light_gray'])
    add_side_accent(slide)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Images
    n_images = len(images_with_captions)
    img_width = 2.6
    spacing = 0.4
    total_width = img_width * n_images + spacing * (n_images - 1)
    start_left = (10 - total_width) / 2
    
    for i, (img_path, caption) in enumerate(images_with_captions):
        left = Inches(start_left + i * (img_width + spacing))
        top = Inches(1.6)
        
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(start_left + i * (img_width + spacing) - 0.1),
            Inches(1.5),
            Inches(img_width + 0.2),
            Inches(5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = RGBColor(220, 220, 230)
        
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width=Inches(img_width))
        
        # Caption
        cap_box = slide.shapes.add_textbox(
            Inches(start_left + i * (img_width + spacing) - 0.1),
            Inches(5.8),
            Inches(img_width + 0.2),
            Inches(0.6)
        )
        ctf = cap_box.text_frame
        cp = ctf.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(14)
        cp.font.bold = True
        cp.font.color.rgb = COLORS['primary']
        cp.alignment = PP_ALIGN.CENTER
    
    return slide


def add_architecture_slide(prs):
    """Add architecture diagram slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['light_gray'])
    add_side_accent(slide)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture & Methodology"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Flow boxes - modern rounded style
    boxes = [
        ("Medical Images\n(Chest X-rays)", 0.4, 2.8, 2, 1.2, COLORS['accent']),
        ("Self-Supervised\nPretraining\n\nJEPA | MAE | Supervised", 3, 2.2, 2.5, 2.4, RGBColor(100, 180, 220)),
        ("ViT-Small\nEncoder\n(384-dim)", 6.2, 2.8, 2, 1.2, RGBColor(130, 200, 160)),
    ]
    
    for text, left, top, width, height, color in boxes:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        
        stf = shape.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = text
        sp.font.size = Pt(12)
        sp.font.bold = True
        sp.font.color.rgb = COLORS['white']
        sp.alignment = PP_ALIGN.CENTER
    
    # Evaluation row
    eval_boxes = [
        ("OOD Detection\nEnergy | Mahalanobis", 0.5, 5, 2.8, 1.3),
        ("Linear Probe\nClassification", 3.6, 5, 2.8, 1.3),
        ("Robustness\nNoise | Blur | Contrast", 6.7, 5, 2.8, 1.3),
    ]
    
    for text, left, top, width, height in eval_boxes:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['primary']
        shape.line.fill.background()
        
        stf = shape.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = text
        sp.font.size = Pt(11)
        sp.font.bold = True
        sp.font.color.rgb = COLORS['white']
        sp.alignment = PP_ALIGN.CENTER
    
    # Arrows
    arrows = [
        (2.4, 3.4, 3, 3.4),
        (5.5, 3.4, 6.2, 3.4),
        (5, 4.6, 5, 5),
    ]
    
    for x1, y1, x2, y2 in arrows:
        line = slide.shapes.add_connector(
            1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        line.line.color.rgb = COLORS['primary_dark']
        line.line.width = Pt(2)
    
    return slide


def add_thank_you_slide(prs):
    """Add a thank you slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['primary_dark'])
    add_corner_decoration(slide)
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), Inches(4),
        Inches(2), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(1))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Questions & Discussion"
    sp.font.size = Pt(24)
    sp.font.color.rgb = COLORS['accent2']
    sp.alignment = PP_ALIGN.CENTER
    
    return slide


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: Title ==========
    add_title_slide(
        prs,
        "Self-Supervised Learning for\nMedical Imaging OOD Detection",
        "Comparing JEPA, MAE, and Supervised Learning for Chest X-ray Analysis"
    )
    
    # ========== SLIDE 2: Problem & Motivation ==========
    add_content_slide(prs, "Problem Statement & Motivation", [
        "Medical AI models encounter out-of-distribution (OOD) data in clinical settings",
        "   Different scanners, institutions, and patient populations",
        "Reliable OOD detection is essential for safe AI deployment",
        "   Models should know when predictions may be unreliable",
        "Self-supervised learning (SSL) shows promise for robust representations",
        "   JEPA: Predicts latent embeddings (semantic features)",
        "   MAE: Pixel reconstruction (intensity-invariant features)",
        "Goal: Compare SSL methods for medical OOD detection"
    ])
    
    # ========== SLIDE 3: Methodology & Architecture ==========
    add_architecture_slide(prs)
    
    # ========== SLIDE 4: Methods Deep Dive ==========
    add_content_slide(prs, "Self-Supervised Learning Methods", [
        "JEPA (Joint Embedding Predictive Architecture)",
        "   Predicts target patch embeddings from context (85-100% masking)",
        "   Context + Target encoders with EMA update",
        "MAE (Masked Autoencoder)",
        "   Reconstructs masked pixels (75% masking ratio)",
        "   Asymmetric encoder-decoder architecture",
        "OOD Detection Scoring",
        "   Energy Score: E(x) = -log Σ exp(f_c(x))",
        "   Mahalanobis Distance: Models embedding distribution"
    ], two_column=False)
    
    # ========== SLIDE 5: Datasets ==========
    add_content_slide(prs, "Datasets & Experimental Setup", [
        "In-Distribution (Training)",
        "   TB Chest Radiography Database: 4,200 images",
        "   Binary: Normal (700) vs Tuberculosis (3,500)",
        "Out-of-Distribution (Evaluation Only)",
        "   Montgomery TB: 138 images (US hospital)",
        "   Shenzhen TB: 662 images (China hospital)",
        "Training Configuration",
        "   ViT-Small: 21.6M parameters, 384-dim embeddings",
        "   50 epochs, batch size 32, cosine LR decay"
    ], two_column=False)
    
    # ========== SLIDE 6: Libraries & Tools ==========
    add_content_slide(prs, "Libraries & Technologies", [
        "Deep Learning: PyTorch 2.0, torchvision",
        "Architecture: Vision Transformer (ViT-Small)",
        "SSL Framework: I-JEPA (Meta AI), MAE",
        "Visualization: Matplotlib, Plotly, PyVista, t-SNE",
        "Metrics: scikit-learn (AUROC, covariance estimation)",
        "Config & Utils: YAML, NumPy, Pandas, PIL"
    ], two_column=False)
    
    # ========== SLIDE 7: Embedding Visualizations ==========
    images = [
        (os.path.join(EMBEDDING_DIR, "embedding_3d_jepa.png"), "JEPA"),
        (os.path.join(EMBEDDING_DIR, "embedding_3d_mae.png"), "MAE"),
        (os.path.join(EMBEDDING_DIR, "embedding_3d_supervised.png"), "Supervised")
    ]
    add_multi_image_slide(prs, "Embedding Space Visualization (t-SNE)", images)
    
    # ========== SLIDE 8: Future Work ==========
    add_content_slide(prs, "Next Steps & Future Work", [
        "Extend to EchoNet Segmentation",
        "   JEPA pretraining on echocardiogram videos",
        "   Left ventricular segmentation task",
        "   Transfer to pediatric echo (OOD evaluation)",
        "Model & Evaluation Improvements",
        "   Larger ViT architectures (ViT-Base, ViT-Large)",
        "   Video-level JEPA for temporal modeling",
        "   Additional domains: CT, MRI",
        "Clinical Validation",
        "   Multi-site deployment testing"
    ], two_column=False)
    
    # ========== SLIDE 9: Thank You ==========
    add_thank_you_slide(prs)
    
    # Save presentation
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    return OUTPUT_PATH


if __name__ == "__main__":
    create_presentation()
