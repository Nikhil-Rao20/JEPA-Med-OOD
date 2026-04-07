#!/usr/bin/env python3
"""
Generate Modern PowerPoint Presentation with Results for Med_JEPA_ODD Project
Self-Supervised Learning for Medical Imaging OOD Detection
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Get the project root directory
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(PROJECT_ROOT, "Med_JEPA_OOD_Presentation_Final.pptx")

# Path to images
EMBEDDING_DIR = os.path.join(PROJECT_ROOT, "experiments", "cxr_jepa_pilot", "3d_visualization")
ANALYSIS_DIR = os.path.join(PROJECT_ROOT, "experiments", "cxr_jepa_pilot", "embedding_analysis")
ROBUSTNESS_DIR = os.path.join(PROJECT_ROOT, "experiments", "cxr_jepa_pilot", "robustness_ablation")

# Modern Color Palette - Deep Blue/Teal theme
COLORS = {
    'primary_dark': RGBColor(15, 32, 65),
    'primary': RGBColor(26, 54, 93),
    'accent': RGBColor(0, 180, 216),
    'accent2': RGBColor(144, 224, 239),
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(240, 248, 255),
    'text_dark': RGBColor(30, 30, 50),
    'green': RGBColor(46, 204, 113),
    'red': RGBColor(231, 76, 60),
    'orange': RGBColor(243, 156, 18),
}


def add_solid_background(slide, color):
    """Add a solid color background."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_side_accent(slide):
    """Add a vertical accent bar on the left side."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.15), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()


def add_corner_decoration(slide):
    """Add decorative elements in corners."""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8.5), Inches(-0.5),
        Inches(2), Inches(2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['accent']
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-0.8), Inches(6),
        Inches(2.5), Inches(2.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = COLORS['accent2']
    circle2.fill.fore_color.brightness = 0.4
    circle2.line.fill.background()


def add_title_slide(prs, title, subtitle=""):
    """Add a modern title slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['primary_dark'])
    add_corner_decoration(slide)
    
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(3.3),
        Inches(2), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(1.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1))
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = COLORS['accent2']
        sp.alignment = PP_ALIGN.LEFT
    
    return slide


def add_content_slide(prs, title, bullet_points):
    """Add a modern content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['light_gray'])
    add_side_accent(slide)
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = ctf.paragraphs[0]
        else:
            p = ctf.add_paragraph()
        
        if point.startswith("   "):
            p.text = "    → " + point.strip()
            p.font.size = Pt(16)
        else:
            p.text = "● " + point
            p.font.size = Pt(18)
        
        p.font.color.rgb = COLORS['text_dark']
        p.space_after = Pt(8)
    
    return slide


def add_table_slide(prs, title, headers, data, highlight_col=None, highlight_best=True):
    """Add a slide with a data table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['light_gray'])
    add_side_accent(slide)
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Create table
    rows = len(data) + 1
    cols = len(headers)
    
    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(9)
    height = Inches(0.5 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Style header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = COLORS['white']
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Find best values per column for highlighting
    best_values = {}
    if highlight_best and len(data) > 0:
        for col_idx in range(1, len(headers)):
            try:
                values = [float(row[col_idx].replace('%', '')) for row in data if row[col_idx] not in ['', '-']]
                if values:
                    best_values[col_idx] = max(values)
            except:
                pass
    
    # Fill data rows
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 240, 250)
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(13)
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Highlight best values
            if col_idx in best_values:
                try:
                    if float(cell_text.replace('%', '')) == best_values[col_idx]:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = COLORS['green']
                except:
                    pass
            
            # First column (method names) - left align and bold
            if col_idx == 0:
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.font.bold = True
                paragraph.font.color.rgb = COLORS['primary_dark']
    
    return slide


def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['light_gray'])
    add_side_accent(slide)
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    if os.path.exists(image_path):
        img_left = Inches(1.5)
        img_top = Inches(1.5)
        img_width = Inches(7)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
    
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        ctf = cap_box.text_frame
        cp = ctf.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(12)
        cp.font.italic = True
        cp.font.color.rgb = COLORS['text_dark']
        cp.alignment = PP_ALIGN.CENTER
    
    return slide


def add_multi_image_slide(prs, title, images_with_captions):
    """Add a slide with multiple images."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['light_gray'])
    add_side_accent(slide)
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    n_images = len(images_with_captions)
    img_width = 2.6
    spacing = 0.4
    total_width = img_width * n_images + spacing * (n_images - 1)
    start_left = (10 - total_width) / 2
    
    for i, (img_path, caption) in enumerate(images_with_captions):
        left = Inches(start_left + i * (img_width + spacing))
        top = Inches(1.6)
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(start_left + i * (img_width + spacing) - 0.1),
            Inches(1.5),
            Inches(img_width + 0.2),
            Inches(5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = RGBColor(220, 220, 230)
        
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width=Inches(img_width))
        
        cap_box = slide.shapes.add_textbox(
            Inches(start_left + i * (img_width + spacing) - 0.1),
            Inches(5.8),
            Inches(img_width + 0.2),
            Inches(0.6)
        )
        ctf = cap_box.text_frame
        cp = ctf.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(14)
        cp.font.bold = True
        cp.font.color.rgb = COLORS['primary']
        cp.alignment = PP_ALIGN.CENTER
    
    return slide


def add_key_findings_slide(prs):
    """Add a key findings slide with icons."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['light_gray'])
    add_side_accent(slide)
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Findings"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    findings = [
        ("SSL > Supervised for OOD", "Self-supervised methods outperform supervised by 15-25% on energy-based OOD detection", COLORS['green']),
        ("Mahalanobis Works for All", "95-99% AUROC across all methods - distribution-aware scoring compensates for representations", COLORS['accent']),
        ("MAE = Best Robustness", "MAE maintains 98% AUC under contrast corruption where others drop to 67-76%", COLORS['green']),
        ("MAE = Best Transfer", "MAE outperforms on cross-dataset generalization by 10-26 percentage points", COLORS['green']),
    ]
    
    for i, (title, desc, color) in enumerate(findings):
        top = 1.5 + i * 1.4
        
        # Icon circle
        icon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5), Inches(top),
            Inches(0.6), Inches(0.6)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()
        
        # Number in icon
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(top + 0.1), Inches(0.6), Inches(0.4))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = COLORS['white']
        np.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1.3), Inches(top), Inches(8), Inches(0.4))
        ttf = title_box.text_frame
        tp = ttf.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(18)
        tp.font.bold = True
        tp.font.color.rgb = COLORS['primary_dark']
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.3), Inches(top + 0.4), Inches(8), Inches(0.8))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(14)
        dp.font.color.rgb = COLORS['text_dark']
    
    return slide


def add_thank_you_slide(prs):
    """Add a thank you slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_solid_background(slide, COLORS['primary_dark'])
    add_corner_decoration(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), Inches(4),
        Inches(2), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(1))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Questions & Discussion"
    sp.font.size = Pt(24)
    sp.font.color.rgb = COLORS['accent2']
    sp.alignment = PP_ALIGN.CENTER
    
    return slide


def create_presentation():
    """Create the full presentation with results."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: Title ==========
    add_title_slide(
        prs,
        "Self-Supervised Learning for\nMedical Imaging OOD Detection",
        "Comparing JEPA, MAE, and Supervised Learning for Chest X-ray Analysis"
    )
    
    # ========== SLIDE 2: Problem & Motivation ==========
    add_content_slide(prs, "Problem Statement & Motivation", [
        "Medical AI models encounter out-of-distribution (OOD) data in clinical settings",
        "   Different scanners, institutions, and patient populations",
        "Reliable OOD detection is essential for safe AI deployment",
        "   Models should know when predictions may be unreliable",
        "Self-supervised learning (SSL) shows promise for robust representations",
        "   JEPA: Predicts latent embeddings (semantic features)",
        "   MAE: Pixel reconstruction (intensity-invariant features)",
        "Goal: Compare SSL methods for medical OOD detection"
    ])
    
    # ========== SLIDE 3: Methodology ==========
    add_content_slide(prs, "Methodology & Experimental Setup", [
        "Architecture: ViT-Small (21.6M params, 384-dim embeddings)",
        "In-Distribution: TB Chest Radiography Database (4,200 images)",
        "   80/20 train/validation split (3,360 / 840 samples)",
        "Out-of-Distribution (Evaluation Only):",
        "   Montgomery TB: 138 images (US hospital)",
        "   Shenzhen TB: 662 images (China hospital)",
        "Training: 50 epochs, batch size 32, AdamW optimizer",
        "OOD Scoring: Energy Score & Mahalanobis Distance"
    ])
    
    # ========== SLIDE 4: ID Classification Results ==========
    add_table_slide(
        prs,
        "Results: In-Distribution Classification",
        ["Method", "AUROC", "Accuracy", "ECE", "NLL"],
        [
            ["JEPA", "93.9%", "91.7%", "0.029", "0.206"],
            ["MAE", "99.3%", "97.1%", "0.013", "0.085"],
            ["Supervised", "99.6%", "98.0%", "0.007", "0.064"],
        ]
    )
    
    # ========== SLIDE 5: OOD Detection Results ==========
    add_table_slide(
        prs,
        "Results: OOD Detection (Energy Score)",
        ["Method", "Mont. AUROC", "Mont. FPR@95", "Shen. AUROC", "Shen. FPR@95"],
        [
            ["JEPA", "77.1%", "54.0%", "71.8%", "63.3%"],
            ["MAE", "76.9%", "57.7%", "86.6%", "35.8%"],
            ["Supervised", "60.7%", "81.4%", "58.2%", "93.1%"],
        ]
    )
    
    # ========== SLIDE 6: Mahalanobis OOD Detection ==========
    add_table_slide(
        prs,
        "Results: OOD Detection (Mahalanobis Distance)",
        ["Method", "Mont. AUROC", "Mont. FPR@95", "Shen. AUROC", "Shen. FPR@95"],
        [
            ["JEPA", "95.6%", "18.6%", "98.5%", "7.7%"],
            ["MAE", "99.4%", "3.0%", "98.4%", "8.5%"],
            ["Supervised", "99.1%", "2.4%", "99.7%", "0.0%"],
        ]
    )
    
    # ========== SLIDE 7: Robustness Results ==========
    add_table_slide(
        prs,
        "Results: Robustness to Image Corruptions (AUC at Severity 5)",
        ["Method", "Clean", "Noise", "Blur", "Contrast"],
        [
            ["JEPA", "99.0%", "98.3%", "98.1%", "67.2%"],
            ["MAE", "99.2%", "98.4%", "99.2%", "98.3%"],
            ["Supervised", "99.6%", "99.2%", "99.4%", "75.8%"],
        ]
    )
    
    # ========== SLIDE 8: Embedding Visualizations ==========
    images = [
        (os.path.join(EMBEDDING_DIR, "embedding_3d_jepa.png"), "JEPA"),
        (os.path.join(EMBEDDING_DIR, "embedding_3d_mae.png"), "MAE"),
        (os.path.join(EMBEDDING_DIR, "embedding_3d_supervised.png"), "Supervised")
    ]
    add_multi_image_slide(prs, "Embedding Space Visualization (t-SNE)", images)
    
    # ========== SLIDE 9: Key Findings ==========
    add_key_findings_slide(prs)
    
    # ========== SLIDE 10: Future Work ==========
    add_content_slide(prs, "Future Work: EchoNet Segmentation", [
        "Extend framework to echocardiogram segmentation",
        "   EchoNet-Dynamic: In-distribution training (adult echo)",
        "   EchoNet-Pediatric: OOD evaluation (pediatric echo)",
        "Preliminary Results:",
        "   MAE achieves 91.2% Dice on ID (vs 84.7% JEPA, 80.4% Supervised)",
        "   MAE transfers best to Pediatric A4C: 83.4% Dice",
        "Additional Directions:",
        "   Scale to larger ViT architectures and foundation models",
        "   Temporal JEPA variants for video understanding",
        "   Multi-site clinical validation"
    ])
    
    # ========== SLIDE 11: Thank You ==========
    add_thank_you_slide(prs)
    
    # Save presentation
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    return OUTPUT_PATH


if __name__ == "__main__":
    create_presentation()
